#!/opt/hermes/.venv/bin/python
from __future__ import annotations

import argparse
import csv
import html
import json
import shutil
import subprocess
import tempfile
import zipfile
from pathlib import Path
from typing import Any, Iterable


OFFICE_TARGETS = {".doc": "docx", ".xls": "xlsx", ".ppt": "pptx"}
SUPPORTED_OUTPUTS = {".docx", ".pdf", ".xlsx", ".csv", ".pptx"}


def emit(data: dict[str, Any]) -> None:
    print(json.dumps(data, ensure_ascii=False, indent=2))


def run(command: list[str], *, timeout: int = 180) -> subprocess.CompletedProcess[str]:
    return subprocess.run(command, check=True, capture_output=True, text=True, timeout=timeout)


def load_spec(path: Path) -> dict[str, Any]:
    value = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(value, dict):
        raise ValueError("规范文件必须是 JSON 对象")
    return value


def modern_source(path: Path, temporary: Path) -> Path:
    target_format = OFFICE_TARGETS.get(path.suffix.lower())
    if not target_format:
        return path
    run(["soffice", "--headless", "--convert-to", target_format, "--outdir", str(temporary), str(path)])
    result = temporary / f"{path.stem}.{target_format}"
    if not result.is_file():
        matches = list(temporary.glob(f"*.{target_format}"))
        if not matches:
            raise RuntimeError(f"LibreOffice 未能转换 {path.name}")
        result = matches[0]
    return result


def markdown_table(rows: Iterable[Iterable[Any]]) -> str:
    values = [[str(value if value is not None else "").replace("|", "\\|").replace("\n", "<br>") for value in row] for row in rows]
    if not values:
        return ""
    width = max(len(row) for row in values)
    values = [row + [""] * (width - len(row)) for row in values]
    return "\n".join([
        "| " + " | ".join(values[0]) + " |",
        "| " + " | ".join(["---"] * width) + " |",
        *("| " + " | ".join(row) + " |" for row in values[1:]),
    ])


def inspect_docx(path: Path) -> tuple[str, dict[str, Any]]:
    from docx import Document

    document = Document(path)
    parts: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = str(paragraph.style.name or "")
        if style.lower().startswith("heading"):
            try:
                level = min(6, max(1, int(style.split()[-1])))
            except ValueError:
                level = 2
            parts.append(f"{'#' * level} {text}")
        else:
            parts.append(text)
    for index, table in enumerate(document.tables, 1):
        rendered = markdown_table([[cell.text for cell in row.cells] for row in table.rows])
        if rendered:
            parts.extend([f"## 表格 {index}", rendered])
    return "\n\n".join(parts), {
        "paragraphs": len(document.paragraphs),
        "tables": len(document.tables),
        "sections": len(document.sections),
    }


def inspect_pdf(path: Path, ocr: str) -> tuple[str, dict[str, Any]]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    page_text = [(page.extract_text() or "").strip() for page in reader.pages]
    direct_chars = sum(len(value) for value in page_text)
    ocr_indices = set(range(len(page_text))) if ocr == "always" else {
        index for index, value in enumerate(page_text) if ocr == "auto" and len(value) < 20
    }
    ocr_pages = 0
    if ocr_indices:
        with tempfile.TemporaryDirectory(prefix="document-ocr-") as directory:
            prefix = Path(directory) / "page"
            run(["pdftoppm", "-png", "-r", "220", str(path), str(prefix)], timeout=600)
            images = sorted(Path(directory).glob("page-*.png"))
            if not images:
                raise RuntimeError("PDF OCR 渲染未生成页面")
            for index, image_path in enumerate(images):
                if index not in ocr_indices:
                    continue
                result = run(["tesseract", str(image_path), "stdout", "-l", "chi_sim+eng", "--psm", "3"], timeout=180)
                page_text[index] = result.stdout.strip()
                ocr_pages += 1
    text = "\n\n".join(f"## 第 {index} 页\n\n{value}" for index, value in enumerate(page_text, 1) if value)
    return text, {
        "pages": len(reader.pages),
        "directTextCharacters": direct_chars,
        "ocrUsed": bool(ocr_indices),
        "ocrPages": ocr_pages,
    }


def inspect_xlsx(path: Path, max_rows: int) -> tuple[str, dict[str, Any]]:
    from openpyxl import load_workbook

    workbook = load_workbook(path, data_only=False, read_only=True)
    sections: list[str] = []
    metadata: list[dict[str, Any]] = []
    for sheet in workbook.worksheets:
        rows: list[tuple[Any, ...]] = []
        total_rows = 0
        total_columns = 0
        for index, row in enumerate(sheet.iter_rows(values_only=True), 1):
            if not any(value is not None for value in row):
                continue
            total_rows = index
            total_columns = max(total_columns, max((column for column, value in enumerate(row, 1) if value is not None), default=0))
            if max_rows <= 0 or len(rows) < max_rows:
                rows.append(row[:total_columns or len(row)])
        sections.extend([f"# 工作表：{sheet.title}", markdown_table(rows)])
        metadata.append({"name": sheet.title, "rows": total_rows, "columns": total_columns, "truncated": max_rows > 0 and len(rows) < total_rows})
    workbook.close()
    return "\n\n".join(section for section in sections if section), {"sheets": metadata}


def inspect_xls(path: Path, max_rows: int) -> tuple[str, dict[str, Any]]:
    import xlrd

    workbook = xlrd.open_workbook(path)
    sections: list[str] = []
    metadata: list[dict[str, Any]] = []
    for sheet in workbook.sheets():
        limit = sheet.nrows if max_rows <= 0 else min(sheet.nrows, max_rows)
        rows = [[sheet.cell_value(row, column) for column in range(sheet.ncols)] for row in range(limit)]
        sections.extend([f"# 工作表：{sheet.name}", markdown_table(rows)])
        metadata.append({"name": sheet.name, "rows": sheet.nrows, "columns": sheet.ncols, "truncated": limit < sheet.nrows})
    return "\n\n".join(section for section in sections if section), {"sheets": metadata}


def inspect_csv(path: Path, max_rows: int) -> tuple[str, dict[str, Any]]:
    raw = path.read_bytes()
    encoding = "utf-8-sig"
    try:
        text = raw.decode(encoding)
    except UnicodeDecodeError:
        encoding = "gb18030"
        text = raw.decode(encoding)
    sample = text[:8192]
    try:
        dialect = csv.Sniffer().sniff(sample)
    except csv.Error:
        dialect = csv.excel
    rows: list[list[str]] = []
    total = 0
    for row in csv.reader(text.splitlines(), dialect):
        total += 1
        if max_rows <= 0 or len(rows) < max_rows:
            rows.append(row)
    return markdown_table(rows), {"rows": total, "columns": max((len(row) for row in rows), default=0), "encoding": encoding, "truncated": len(rows) < total}


def inspect_pptx(path: Path) -> tuple[str, dict[str, Any]]:
    from pptx import Presentation

    presentation = Presentation(path)
    sections: list[str] = []
    for index, slide in enumerate(presentation.slides, 1):
        lines: list[str] = []
        title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.has_text_frame else ""
        lines.append(f"# 幻灯片 {index}" + (f"：{title}" if title else ""))
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if getattr(shape, "has_text_frame", False):
                value = shape.text.strip()
                if value:
                    lines.append(value)
            if getattr(shape, "has_table", False):
                lines.append(markdown_table([[cell.text for cell in row.cells] for row in shape.table.rows]))
        sections.append("\n\n".join(lines))
    return "\n\n".join(sections), {"slides": len(presentation.slides), "width": presentation.slide_width, "height": presentation.slide_height}


def inspect_file(path: Path, ocr: str, max_rows: int) -> tuple[str, dict[str, Any]]:
    if not path.is_file():
        raise FileNotFoundError(path)
    suffix = path.suffix.lower()
    with tempfile.TemporaryDirectory(prefix="document-convert-") as directory:
        source = modern_source(path, Path(directory))
        suffix = source.suffix.lower()
        if suffix == ".docx":
            return inspect_docx(source)
        if suffix == ".pdf":
            return inspect_pdf(source, ocr)
        if suffix == ".xlsx":
            return inspect_xlsx(source, max_rows)
        if suffix == ".xls":
            return inspect_xls(source, max_rows)
        if suffix in {".csv", ".tsv"}:
            return inspect_csv(source, max_rows)
        if suffix == ".pptx":
            return inspect_pptx(source)
    raise ValueError(f"不支持的文件格式：{path.suffix}")


def blocks(spec: dict[str, Any]) -> list[dict[str, Any]]:
    value = spec.get("blocks", spec.get("content", []))
    if not isinstance(value, list):
        raise ValueError("blocks 必须是数组")
    return [item for item in value if isinstance(item, dict)]


def build_docx(spec: dict[str, Any], output: Path) -> None:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Cm, Pt

    document = Document()
    section = document.sections[0]
    section.top_margin = Cm(float(spec.get("marginCm", 2.2)))
    section.bottom_margin = Cm(float(spec.get("marginCm", 2.2)))
    normal = document.styles["Normal"]
    normal.font.name = str(spec.get("font", "Noto Sans CJK SC"))
    normal.font.size = Pt(float(spec.get("fontSize", 10.5)))
    title = str(spec.get("title", "")).strip()
    if title:
        document.add_heading(title, 0)
    for block in blocks(spec):
        kind = str(block.get("type", "paragraph"))
        if kind == "heading":
            document.add_heading(str(block.get("text", "")), level=min(6, max(1, int(block.get("level", 1)))))
        elif kind == "paragraph":
            paragraph = document.add_paragraph(str(block.get("text", "")), style=block.get("style"))
            alignment = str(block.get("align", "")).lower()
            if alignment == "center": paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            if alignment == "right": paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        elif kind == "bullets":
            for item in block.get("items", []):
                document.add_paragraph(str(item), style="List Bullet")
        elif kind == "table":
            rows = block.get("rows", [])
            if not rows:
                continue
            width = max(len(row) for row in rows)
            table = document.add_table(rows=len(rows), cols=width)
            table.style = str(block.get("style", "Table Grid"))
            for row_index, row in enumerate(rows):
                for column_index, value in enumerate(row):
                    table.cell(row_index, column_index).text = str(value if value is not None else "")
        elif kind == "pageBreak":
            document.add_page_break()
        elif kind == "image":
            document.add_picture(str(block["path"]), width=Cm(float(block.get("widthCm", 15))))
    document.core_properties.title = title
    document.core_properties.author = str(spec.get("author", "也算是学 Agent"))
    output.parent.mkdir(parents=True, exist_ok=True)
    document.save(output)


def cjk_font() -> str:
    candidates = [
        "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSerifCJK-Regular.ttc",
    ]
    for candidate in candidates:
        if Path(candidate).is_file():
            return candidate
    raise RuntimeError("未找到可嵌入 PDF 的中文字体")


def build_pdf(spec: dict[str, Any], output: Path) -> None:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

    pdfmetrics.registerFont(TTFont("NotoCJK", cjk_font(), subfontIndex=0))
    styles = getSampleStyleSheet()
    body = ParagraphStyle("CJKBody", parent=styles["BodyText"], fontName="NotoCJK", fontSize=10.5, leading=17, spaceAfter=7)
    headings = {level: ParagraphStyle(f"CJKH{level}", parent=styles[f"Heading{min(level, 3)}"], fontName="NotoCJK", fontSize=19 - level * 2, leading=25 - level * 2, spaceBefore=8, spaceAfter=8) for level in range(1, 7)}
    title_style = ParagraphStyle("CJKTitle", parent=styles["Title"], fontName="NotoCJK", fontSize=22, leading=30, alignment=TA_CENTER, spaceAfter=16)
    output.parent.mkdir(parents=True, exist_ok=True)
    document = SimpleDocTemplate(str(output), pagesize=A4, leftMargin=2 * cm, rightMargin=2 * cm, topMargin=2 * cm, bottomMargin=2 * cm, title=str(spec.get("title", "")), author=str(spec.get("author", "也算是学 Agent")))
    story: list[Any] = []
    if spec.get("title"):
        story.append(Paragraph(html.escape(str(spec["title"])), title_style))
    for block in blocks(spec):
        kind = str(block.get("type", "paragraph"))
        if kind == "heading":
            level = min(6, max(1, int(block.get("level", 1))))
            story.append(Paragraph(html.escape(str(block.get("text", ""))).replace("\n", "<br/>"), headings[level]))
        elif kind == "paragraph":
            story.append(Paragraph(html.escape(str(block.get("text", ""))).replace("\n", "<br/>"), body))
        elif kind == "bullets":
            for item in block.get("items", []):
                story.append(Paragraph("• " + html.escape(str(item)), body))
        elif kind == "table":
            rows = [[Paragraph(html.escape(str(value if value is not None else "")), body) for value in row] for row in block.get("rows", [])]
            if rows:
                table = Table(rows, repeatRows=1, hAlign="LEFT")
                table.setStyle(TableStyle([
                    ("FONTNAME", (0, 0), (-1, -1), "NotoCJK"),
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#ECEEEF")),
                    ("GRID", (0, 0), (-1, -1), 0.45, colors.HexColor("#BFC4C6")),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 5),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 5),
                ]))
                story.extend([table, Spacer(1, 8)])
        elif kind == "pageBreak":
            story.append(PageBreak())
    document.build(story)


def build_xlsx(spec: dict[str, Any], output: Path) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    from openpyxl.utils import get_column_letter

    workbook = Workbook()
    workbook.remove(workbook.active)
    sheets = spec.get("sheets", [])
    if not sheets:
        sheets = [{"name": "Sheet1", "rows": spec.get("rows", [])}]
    for sheet_spec in sheets:
        sheet = workbook.create_sheet(str(sheet_spec.get("name", "Sheet"))[:31])
        rows = sheet_spec.get("rows", [])
        for row in rows:
            sheet.append(list(row))
        header_rows = max(0, int(sheet_spec.get("headerRows", 1 if rows else 0)))
        thin = Side(style="thin", color="D7DADC")
        for row in sheet.iter_rows():
            for cell in row:
                cell.alignment = Alignment(vertical="top", wrap_text=True)
                cell.border = Border(bottom=thin)
                if cell.row <= header_rows:
                    cell.font = Font(bold=True, color="303336")
                    cell.fill = PatternFill("solid", fgColor="ECEEEF")
        widths = sheet_spec.get("columnWidths", {})
        for index in range(1, sheet.max_column + 1):
            explicit = widths.get(get_column_letter(index)) if isinstance(widths, dict) else None
            measured = max((len(str(sheet.cell(row, index).value or "")) for row in range(1, min(sheet.max_row, 500) + 1)), default=8)
            sheet.column_dimensions[get_column_letter(index)].width = float(explicit or min(48, max(10, measured + 2)))
        if sheet_spec.get("freezePanes"):
            sheet.freeze_panes = str(sheet_spec["freezePanes"])
        elif header_rows:
            sheet.freeze_panes = f"A{header_rows + 1}"
        if sheet_spec.get("autoFilter", bool(header_rows and sheet.max_column)) and sheet.max_row:
            sheet.auto_filter.ref = sheet.dimensions
        for address, value in (sheet_spec.get("numberFormats") or {}).items():
            sheet[str(address)].number_format = str(value)
    output.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(output)


def build_csv(spec: dict[str, Any], output: Path) -> None:
    output.parent.mkdir(parents=True, exist_ok=True)
    with output.open("w", encoding="utf-8-sig", newline="") as stream:
        writer = csv.writer(stream)
        writer.writerows(spec.get("rows", []))


def build_pptx(spec: dict[str, Any], output: Path) -> None:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slides = spec.get("slides", [])
    if not slides:
        raise ValueError("slides 不能为空")
    for index, slide_spec in enumerate(slides):
        layout = presentation.slide_layouts[0 if index == 0 and slide_spec.get("subtitle") is not None else 5]
        slide = presentation.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = str(slide_spec.get("title", ""))
        if index == 0 and slide_spec.get("subtitle") is not None and len(slide.placeholders) > 1:
            slide.placeholders[1].text = str(slide_spec.get("subtitle", ""))
        y = 1.55
        bullets = slide_spec.get("bullets", [])
        if bullets:
            box = slide.shapes.add_textbox(Inches(0.8), Inches(y), Inches(11.7), Inches(4.9))
            frame = box.text_frame
            frame.clear()
            for bullet_index, item in enumerate(bullets):
                paragraph = frame.paragraphs[0] if bullet_index == 0 else frame.add_paragraph()
                if isinstance(item, dict):
                    paragraph.text = str(item.get("text", "")); paragraph.level = int(item.get("level", 0))
                else:
                    paragraph.text = str(item)
                paragraph.font.size = Pt(24 if paragraph.level == 0 else 18)
                paragraph.space_after = Pt(10)
            y = 5.75
        if slide_spec.get("body"):
            box = slide.shapes.add_textbox(Inches(0.85), Inches(y), Inches(11.6), Inches(4.8))
            frame = box.text_frame
            frame.word_wrap = True
            frame.text = str(slide_spec["body"])
            frame.paragraphs[0].font.size = Pt(22)
        rows = slide_spec.get("table")
        if rows:
            row_count, column_count = len(rows), max(len(row) for row in rows)
            table = slide.shapes.add_table(row_count, column_count, Inches(0.65), Inches(1.55), Inches(12.0), Inches(4.9)).table
            for row_index, row in enumerate(rows):
                for column_index, value in enumerate(row):
                    cell = table.cell(row_index, column_index)
                    cell.text = str(value if value is not None else "")
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(13)
                        if row_index == 0: paragraph.font.bold = True
                        paragraph.alignment = PP_ALIGN.LEFT
                    if row_index == 0:
                        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(236, 238, 239)
        image = slide_spec.get("image")
        if image:
            slide.shapes.add_picture(str(image), Inches(float(slide_spec.get("imageX", 7.0))), Inches(float(slide_spec.get("imageY", 1.5))), width=Inches(float(slide_spec.get("imageWidth", 5.4))))
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)


def build(spec: dict[str, Any]) -> Path:
    output = Path(str(spec.get("output", ""))).expanduser()
    if output.suffix.lower() not in SUPPORTED_OUTPUTS:
        raise ValueError("output 必须是 .docx、.pdf、.xlsx、.csv 或 .pptx")
    builders = {".docx": build_docx, ".pdf": build_pdf, ".xlsx": build_xlsx, ".csv": build_csv, ".pptx": build_pptx}
    builders[output.suffix.lower()](spec, output)
    return output


def replace_text(value: str, replacements: list[dict[str, Any]]) -> str:
    for replacement in replacements:
        value = value.replace(str(replacement.get("old", "")), str(replacement.get("new", "")))
    return value


def edit(spec: dict[str, Any]) -> Path:
    source = Path(str(spec.get("source", ""))).expanduser()
    output = Path(str(spec.get("output", ""))).expanduser()
    if not source.is_file() or output.suffix.lower() not in SUPPORTED_OUTPUTS:
        raise ValueError("source 不存在或 output 格式不受支持")
    replacements = [item for item in spec.get("replacements", []) if isinstance(item, dict)]
    suffix = source.suffix.lower()
    output.parent.mkdir(parents=True, exist_ok=True)
    if suffix == ".docx" and output.suffix.lower() == ".docx":
        from docx import Document
        document = Document(source)
        for paragraph in document.paragraphs:
            if replacements and any(str(item.get("old", "")) in paragraph.text for item in replacements):
                paragraph.text = replace_text(paragraph.text, replacements)
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    cell.text = replace_text(cell.text, replacements)
        document.save(output)
    elif suffix == ".xlsx" and output.suffix.lower() == ".xlsx":
        from openpyxl import load_workbook
        workbook = load_workbook(source)
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str): cell.value = replace_text(cell.value, replacements)
        for change in spec.get("setCells", []):
            workbook[str(change["sheet"])][str(change["cell"])] = change.get("value")
        workbook.save(output)
    elif suffix == ".csv" and output.suffix.lower() == ".csv":
        text = source.read_text(encoding="utf-8-sig")
        output.write_text(replace_text(text, replacements), encoding="utf-8-sig")
    elif suffix == ".pptx" and output.suffix.lower() == ".pptx":
        from pptx import Presentation
        presentation = Presentation(source)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False): shape.text = replace_text(shape.text, replacements)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        for cell in row.cells: cell.text = replace_text(cell.text, replacements)
        presentation.save(output)
    else:
        raise ValueError("edit 目前支持同格式的 DOCX、XLSX、CSV、PPTX；旧格式请先 convert，PDF 请按提取内容重建")
    return output


def render_and_check(path: Path) -> dict[str, Any]:
    from PIL import Image, ImageStat

    with tempfile.TemporaryDirectory(prefix="document-verify-") as directory:
        root = Path(directory)
        pdf = path
        if path.suffix.lower() != ".pdf":
            run(["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(root), str(path)], timeout=300)
            pdf = root / f"{path.stem}.pdf"
            if not pdf.is_file():
                matches = list(root.glob("*.pdf"))
                if not matches:
                    raise RuntimeError("LibreOffice 渲染未生成 PDF")
                pdf = matches[0]
        prefix = root / "render"
        run(["pdftoppm", "-png", "-r", "120", str(pdf), str(prefix)], timeout=300)
        images = sorted(root.glob("render-*.png"))
        if not images:
            raise RuntimeError("页面渲染未生成图片")
        nonblank = 0
        for image_path in images:
            with Image.open(image_path).convert("L") as image:
                extrema = ImageStat.Stat(image).extrema[0]
                if extrema[0] < 250:
                    nonblank += 1
        if nonblank != len(images):
            raise RuntimeError("检测到空白渲染页面")
        return {"renderedPages": len(images), "nonblankPages": nonblank}


def verify(path: Path) -> dict[str, Any]:
    if not path.is_file() or path.stat().st_size <= 0:
        raise ValueError("输出文件不存在或为空")
    suffix = path.suffix.lower()
    if suffix in {".docx", ".xlsx", ".pptx"}:
        with zipfile.ZipFile(path) as archive:
            corrupt = archive.testzip()
            if corrupt:
                raise RuntimeError(f"Office ZIP 结构损坏：{corrupt}")
    text, structure = inspect_file(path, "never", 200)
    from markitdown import MarkItDown
    markitdown_text = str(MarkItDown().convert(str(path)).text_content or "").strip()
    if text.strip() and not markitdown_text:
        raise RuntimeError("MarkItDown 二次读取未返回文本")
    result: dict[str, Any] = {
        "ok": True,
        "path": str(path),
        "sizeBytes": path.stat().st_size,
        "textCharacters": len(text),
        "markitdownCharacters": len(markitdown_text),
        "structure": structure,
    }
    result.update(render_and_check(path))
    return result


def convert(source: Path, output: Path) -> Path:
    output.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="document-convert-") as directory:
        run(["soffice", "--headless", "--convert-to", output.suffix.lstrip("."), "--outdir", directory, str(source)], timeout=300)
        generated = Path(directory) / f"{source.stem}{output.suffix}"
        if not generated.is_file():
            matches = list(Path(directory).glob(f"*{output.suffix}"))
            if not matches:
                raise RuntimeError("格式转换失败")
            generated = matches[0]
        shutil.copy2(generated, output)
    return output


def main() -> int:
    parser = argparse.ArgumentParser(description="也算是学 Agent 固定文档处理工具")
    subparsers = parser.add_subparsers(dest="command", required=True)
    inspect_parser = subparsers.add_parser("inspect")
    inspect_parser.add_argument("path", type=Path)
    inspect_parser.add_argument("--ocr", choices=["auto", "always", "never"], default="auto")
    inspect_parser.add_argument("--max-rows", type=int, default=1000)
    inspect_parser.add_argument("--output", type=Path)
    for command in ["build", "edit"]:
        child = subparsers.add_parser(command)
        child.add_argument("spec", type=Path)
    verify_parser = subparsers.add_parser("verify")
    verify_parser.add_argument("path", type=Path)
    convert_parser = subparsers.add_parser("convert")
    convert_parser.add_argument("source", type=Path)
    convert_parser.add_argument("output", type=Path)
    subparsers.add_parser("schema")
    args = parser.parse_args()
    try:
        if args.command == "inspect":
            text, metadata = inspect_file(args.path, args.ocr, args.max_rows)
            if args.output:
                args.output.parent.mkdir(parents=True, exist_ok=True)
                args.output.write_text(text, encoding="utf-8")
            emit({"ok": True, "path": str(args.path), "text": text if not args.output else "", "textOutput": str(args.output) if args.output else "", "metadata": metadata})
        elif args.command == "build":
            output = build(load_spec(args.spec)); emit(verify(output))
        elif args.command == "edit":
            output = edit(load_spec(args.spec)); emit(verify(output))
        elif args.command == "verify":
            emit(verify(args.path))
        elif args.command == "convert":
            output = convert(args.source, args.output); emit(verify(output))
        elif args.command == "schema":
            emit({
                "ok": True,
                "build": {
                    "common": {"output": "/workspace/result.ext", "title": "可选标题"},
                    "docxOrPdf": {"blocks": [
                        {"type": "heading", "level": 1, "text": "标题"},
                        {"type": "paragraph", "text": "正文"},
                        {"type": "bullets", "items": ["项目一", "项目二"]},
                        {"type": "table", "rows": [["表头", "表头"], ["值", "值"]]},
                        {"type": "pageBreak"},
                    ]},
                    "xlsx": {"sheets": [{"name": "Sheet1", "headerRows": 1, "freezePanes": "A2", "rows": [["列名"], ["值或=公式"]], "numberFormats": {"B2": "¥#,##0.00"}}]},
                    "csv": {"rows": [["列名"], ["值"]]},
                    "pptx": {"slides": [{"title": "标题", "subtitle": "副标题"}, {"title": "标题", "bullets": ["要点"]}, {"title": "标题", "table": [["表头"], ["值"]]}]},
                },
                "edit": {"source": "/workspace/source.docx", "output": "/workspace/result.docx", "replacements": [{"old": "原文", "new": "新文"}], "setCells": [{"sheet": "Sheet1", "cell": "B2", "value": 10}]},
                "workflow": ["inspect 源文件", "生成 JSON 规范或脚本", "build/edit", "检查命令返回的 verify 结果", "最终输出 artifact 标记"],
            })
        return 0
    except Exception as error:
        emit({"ok": False, "error": str(error), "type": type(error).__name__})
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
